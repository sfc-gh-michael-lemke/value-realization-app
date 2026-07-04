"""
Value Realization App
Streamlit-in-Snowflake — Cortex COMPLETE · python-pptx · GitHub API · HTML portfolio
"""

import json
import re
import base64
import io
from datetime import date

import requests
import streamlit as st
from snowflake.snowpark.context import get_active_session
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Brand colours ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x2F, 0x4E)
BLUE   = RGBColor(0x29, 0xB5, 0xE8)
LIGHT  = RGBColor(0x00, 0xA1, 0xE0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x44, 0x44, 0x55)

# ─── Snowpark session (cached) ───────────────────────────────────────────────
@st.cache_resource
def get_session():
    try:
        session = get_active_session()
    except Exception:
        # Local dev fallback — create a real session via named connection or env
        from snowflake.snowpark import Session
        import os
        params = {
            "account":       os.getenv("SNOWFLAKE_ACCOUNT",   "sfcogsops-snowhouse_aws_us_west_2"),
            "user":          os.getenv("SNOWFLAKE_USER",      "MLEMKE"),
            "authenticator": os.getenv("SNOWFLAKE_AUTH",      "externalbrowser"),
            "database":      "SALES",
            "schema":        "SALES_ENGINEERING",
            "warehouse":     "SNOWHOUSE",
            "role":          "SALES_ENGINEER",
        }
        session = Session.builder.configs(params).create()
    # Set database/schema context for container runtime (no implicit context)
    try:
        session.sql("USE DATABASE SALES").collect()
        session.sql("USE SCHEMA SALES_ENGINEERING").collect()
        session.sql("USE WAREHOUSE SNOWHOUSE").collect()
    except Exception:
        pass
    return session


# ═══════════════════════════════════════════════════════════════════════════
# LAYER 1 — Cortex extraction
# ═══════════════════════════════════════════════════════════════════════════

EXTRACTION_PROMPT = """Extract the following fields from the use case description below and return ONLY valid JSON (no markdown fencing, no explanation).

Fields:
- customer_name: string (company name, or "Unknown" if not mentioned)
- industry: string (e.g. Financial Services, Retail, Healthcare)
- business_problem: string (2-3 sentence summary of the challenge)
- solution_overview: string (what technologies/products/approaches were used)
- outcomes: array of strings (quantified metrics if available, e.g. "40% faster query performance")
- time_to_value: string (e.g. "3 months", "6 weeks", "Unknown")
- quote: string or null (any testimonial or direct quote from a stakeholder)
- use_case_tags: array of strings chosen from: Data Engineering, AI/ML, Data Sharing, Cost Optimization, Governance, Real-Time Analytics, Application Development, Migration

Use case description:
{description}

Return ONLY the JSON object."""


def extract_fields(description: str) -> dict:
    session = get_session()
    prompt  = EXTRACTION_PROMPT.format(description=description.replace("'", "\\'"))
    sql = f"SELECT SNOWFLAKE.CORTEX.COMPLETE('claude-sonnet-4-5', $$ {prompt} $$) AS result"
    row  = session.sql(sql).collect()[0]["RESULT"]
    # Strip any accidental markdown fencing
    cleaned = re.sub(r"^```[a-z]*\n?", "", row.strip())
    cleaned = re.sub(r"\n?```$", "", cleaned)
    return json.loads(cleaned)


# ═══════════════════════════════════════════════════════════════════════════
# LAYER 2 — PPTX generator
# ═══════════════════════════════════════════════════════════════════════════

def _set_bg(slide, r, g, b):
    from pptx.oxml.ns import qn
    from lxml import etree
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def _add_textbox(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def generate_pptx(fields: dict) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    customer  = fields.get("customer_name", "Customer")
    industry  = fields.get("industry", "")
    problem   = fields.get("business_problem", "")
    solution  = fields.get("solution_overview", "")
    outcomes  = fields.get("outcomes", [])
    ttv       = fields.get("time_to_value", "")
    quote     = fields.get("quote") or ""
    tags      = fields.get("use_case_tags", [])

    # ── Slide 1 — Title ──────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    _set_bg(sl, 0x0A, 0x2F, 0x4E)                         # navy
    _add_rect(sl, 0, 0, 13.33, 0.08, BLUE)                # top accent bar
    _add_rect(sl, 0, 7.42, 13.33, 0.08, BLUE)             # bottom bar

    _add_textbox(sl, "VALUE REALIZATION STORY",
                 1, 1.5, 11, 1, font_size=14, color=LIGHT, bold=True)
    _add_textbox(sl, customer,
                 1, 2.5, 11, 1.5, font_size=44, bold=True, color=WHITE)
    _add_textbox(sl, industry,
                 1, 4.2, 6, 0.8, font_size=20, color=BLUE)
    _add_textbox(sl, f"Time to Value: {ttv}",
                 1, 5.1, 6, 0.6, font_size=16, color=RGBColor(0xBB, 0xCC, 0xDD))
    # Tag pills (simple text row)
    tag_str = "  ·  ".join(tags)
    _add_textbox(sl, tag_str, 1, 5.9, 11, 0.6, font_size=13,
                 color=RGBColor(0x88, 0xBB, 0xDD))
    # Snowflake wordmark placeholder
    _add_textbox(sl, "❄  Snowflake", 10.5, 6.8, 2.5, 0.5,
                 font_size=13, color=BLUE, align=PP_ALIGN.RIGHT)

    # ── Slide 2 — The Challenge ───────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    _set_bg(sl, 0xFF, 0xFF, 0xFF)
    _add_rect(sl, 0, 0, 13.33, 1.1, NAVY)
    _add_textbox(sl, "THE CHALLENGE",
                 0.5, 0.22, 12, 0.7, font_size=24, bold=True, color=WHITE)
    _add_rect(sl, 0.5, 1.4, 0.08, 4.5, BLUE)              # left accent bar
    _add_textbox(sl, problem,
                 0.8, 1.5, 11.8, 4.5, font_size=20,
                 color=NAVY, wrap=True)

    # ── Slide 3 — The Solution ────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    _set_bg(sl, 0xFF, 0xFF, 0xFF)
    _add_rect(sl, 0, 0, 13.33, 1.1, BLUE)
    _add_textbox(sl, "THE SOLUTION",
                 0.5, 0.22, 12, 0.7, font_size=24, bold=True, color=WHITE)
    _add_textbox(sl, solution,
                 0.8, 1.5, 11.8, 3.8, font_size=18, color=NAVY, wrap=True)
    # Tags row
    _add_rect(sl, 0.5, 5.5, 12.3, 0.06, LIGHT)
    _add_textbox(sl, "  ".join([f"[ {t} ]" for t in tags]),
                 0.5, 5.7, 12, 0.7, font_size=13, color=LIGHT)

    # ── Slide 4 — Business Outcomes ───────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    _set_bg(sl, 0x0A, 0x2F, 0x4E)
    _add_rect(sl, 0, 0, 13.33, 1.1, BLUE)
    _add_textbox(sl, "BUSINESS OUTCOMES",
                 0.5, 0.22, 12, 0.7, font_size=24, bold=True, color=WHITE)

    cols  = [0.6, 4.6, 8.6]
    for i, outcome in enumerate(outcomes[:6]):
        col = cols[i % 3]
        row = 1.4 if i < 3 else 4.0
        _add_rect(sl, col, row, 3.6, 2.3, RGBColor(0x12, 0x45, 0x6E))
        _add_rect(sl, col, row, 3.6, 0.06, BLUE)          # top accent
        _add_textbox(sl, outcome, col + 0.1, row + 0.15,
                     3.4, 2.0, font_size=15, color=WHITE, wrap=True)

    # ── Slide 5 — Why Snowflake ───────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    _set_bg(sl, 0xFF, 0xFF, 0xFF)
    _add_rect(sl, 0, 0, 13.33, 1.1, NAVY)
    _add_textbox(sl, "WHY SNOWFLAKE",
                 0.5, 0.22, 12, 0.7, font_size=24, bold=True, color=WHITE)

    bullets = [
        "Single platform — data, AI, and applications in one place",
        "Elastic compute scales instantly to meet any workload",
        "Secure data sharing without copying or moving data",
        "Built-in Cortex AI for intelligent, governed applications",
        f"Time to Value: {ttv}",
    ]
    for idx, b in enumerate(bullets):
        _add_rect(sl, 0.6, 1.4 + idx * 1.0, 0.12, 0.12, BLUE)
        _add_textbox(sl, b, 0.9, 1.3 + idx * 1.0, 11.5, 0.9,
                     font_size=17, color=NAVY, wrap=True)

    # ── Slide 6 — Testimonial (only if quote exists) ──────────────────────────
    if quote.strip():
        sl = prs.slides.add_slide(blank_layout)
        _set_bg(sl, 0x0A, 0x2F, 0x4E)
        _add_rect(sl, 0, 0, 13.33, 1.1, BLUE)
        _add_textbox(sl, "IN THEIR OWN WORDS",
                     0.5, 0.22, 12, 0.7, font_size=24, bold=True, color=WHITE)
        _add_textbox(sl, "\u201c",
                     0.5, 1.3, 1.5, 1.5, font_size=80, color=BLUE)
        _add_textbox(sl, quote,
                     1.5, 2.0, 10.5, 3.5, font_size=22, color=WHITE, wrap=True)
        _add_textbox(sl, f"\u2014 {customer}",
                     1.5, 5.8, 10, 0.6, font_size=16,
                     color=LIGHT, align=PP_ALIGN.RIGHT)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════
# LAYER 3 — GitHub README push
# ═══════════════════════════════════════════════════════════════════════════

README_TEMPLATE = """# {customer_name} — Value Realization Story

> **Industry:** {industry}  
> **Use Case Tags:** {tags}  
> **Time to Value:** {time_to_value}  
> **Generated:** {date}

---

## The Challenge

{business_problem}

---

## The Solution

{solution_overview}

---

## Business Outcomes

{outcomes_md}

---

{quote_block}

---

*Generated by the [Snowflake Value Realization App](https://www.snowflake.com)*
"""


def build_readme(fields: dict) -> str:
    outcomes_md = "\n".join(f"- {o}" for o in fields.get("outcomes", []))
    quote = fields.get("quote") or ""
    quote_block = (
        f'> *"{quote}"*\n>\n> — {fields.get("customer_name", "")}'
        if quote.strip() else ""
    )
    return README_TEMPLATE.format(
        customer_name=fields.get("customer_name", ""),
        industry=fields.get("industry", ""),
        tags=", ".join(fields.get("use_case_tags", [])),
        time_to_value=fields.get("time_to_value", ""),
        date=str(date.today()),
        business_problem=fields.get("business_problem", ""),
        solution_overview=fields.get("solution_overview", ""),
        outcomes_md=outcomes_md,
        quote_block=quote_block,
    )


def push_readme_to_github(token: str, owner: str, repo: str, readme: str) -> dict:
    headers = {
        "Authorization": f"token {token}",
        "Accept": "application/vnd.github+json",
        "X-GitHub-Api-Version": "2022-11-28",
    }
    url = f"https://api.github.com/repos/{owner}/{repo}/contents/README.md"

    # Check if README already exists (need SHA to update)
    existing = requests.get(url, headers=headers, timeout=10)
    sha = existing.json().get("sha") if existing.status_code == 200 else None

    content_b64 = base64.b64encode(readme.encode()).decode()
    payload = {
        "message": "Update value realization story via Snowflake app",
        "content": content_b64,
    }
    if sha:
        payload["sha"] = sha

    resp = requests.put(url, headers=headers, json=payload, timeout=10)
    return resp.json()


def create_github_repo(token: str, owner: str, repo: str, private: bool = False) -> dict:
    """Create a new GitHub repo (user-owned). For org repos, change endpoint."""
    headers = {
        "Authorization": f"token {token}",
        "Accept": "application/vnd.github+json",
    }
    resp = requests.post(
        "https://api.github.com/user/repos",
        headers=headers,
        json={"name": repo, "private": private, "auto_init": False},
        timeout=10,
    )
    return resp.json()


# ═══════════════════════════════════════════════════════════════════════════
# LAYER 4 — HTML portfolio card
# ═══════════════════════════════════════════════════════════════════════════

HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{customer_name} — Value Realization | Snowflake</title>
<style>
  :root {{
    --navy:  #0A2F4E;
    --blue:  #29B5E8;
    --light: #00A1E0;
    --white: #FFFFFF;
    --gray:  #F4F7FA;
    --text:  #1C2B3A;
  }}
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
    background: var(--gray);
    color: var(--text);
    min-height: 100vh;
  }}

  /* ── Header ── */
  .hero {{
    background: linear-gradient(135deg, var(--navy) 0%, #1a4d7a 100%);
    padding: 48px 64px 40px;
    position: relative;
    overflow: hidden;
  }}
  .hero::before {{
    content: "❄";
    position: absolute;
    right: 60px;
    top: 20px;
    font-size: 120px;
    opacity: 0.07;
  }}
  .hero-label {{
    font-size: 11px;
    font-weight: 700;
    letter-spacing: 3px;
    color: var(--blue);
    text-transform: uppercase;
    margin-bottom: 10px;
  }}
  .hero-title {{
    font-size: 42px;
    font-weight: 800;
    color: var(--white);
    line-height: 1.15;
    margin-bottom: 12px;
  }}
  .hero-sub {{
    font-size: 16px;
    color: #8BB8D8;
    margin-bottom: 20px;
  }}
  .tags {{
    display: flex;
    flex-wrap: wrap;
    gap: 8px;
    margin-top: 16px;
  }}
  .tag {{
    background: rgba(41, 181, 232, 0.15);
    border: 1px solid rgba(41, 181, 232, 0.4);
    color: var(--blue);
    padding: 5px 14px;
    border-radius: 20px;
    font-size: 12px;
    font-weight: 600;
    letter-spacing: 0.5px;
  }}

  /* ── Metric bar ── */
  .metric-bar {{
    background: var(--blue);
    padding: 20px 64px;
    display: flex;
    gap: 48px;
    align-items: center;
    flex-wrap: wrap;
  }}
  .metric {{
    color: var(--white);
  }}
  .metric-label {{
    font-size: 10px;
    letter-spacing: 1.5px;
    text-transform: uppercase;
    opacity: 0.8;
  }}
  .metric-value {{
    font-size: 20px;
    font-weight: 700;
    margin-top: 2px;
  }}

  /* ── Body ── */
  .body {{
    max-width: 1100px;
    margin: 40px auto;
    padding: 0 32px;
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: 24px;
  }}
  .card {{
    background: var(--white);
    border-radius: 12px;
    padding: 32px;
    box-shadow: 0 2px 12px rgba(10,47,78,0.07);
    border-top: 4px solid var(--blue);
  }}
  .card.full {{ grid-column: 1 / -1; }}
  .card-icon {{
    font-size: 28px;
    margin-bottom: 12px;
  }}
  .card-title {{
    font-size: 11px;
    font-weight: 700;
    letter-spacing: 2px;
    color: var(--blue);
    text-transform: uppercase;
    margin-bottom: 14px;
  }}
  .card-body {{
    font-size: 15px;
    line-height: 1.7;
    color: #3a4a5a;
  }}
  .outcome-list {{
    list-style: none;
    padding: 0;
  }}
  .outcome-list li {{
    padding: 10px 0;
    border-bottom: 1px solid #eef2f7;
    display: flex;
    align-items: flex-start;
    gap: 10px;
    font-size: 15px;
    color: #3a4a5a;
  }}
  .outcome-list li:last-child {{ border-bottom: none; }}
  .bullet {{
    display: inline-block;
    width: 8px;
    height: 8px;
    background: var(--blue);
    border-radius: 50%;
    margin-top: 6px;
    flex-shrink: 0;
  }}

  /* ── Quote ── */
  .quote-block {{
    background: linear-gradient(135deg, var(--navy) 0%, #1a4d7a 100%);
    border-radius: 12px;
    padding: 36px 40px;
    grid-column: 1 / -1;
    position: relative;
  }}
  .quote-mark {{
    font-size: 72px;
    color: var(--blue);
    line-height: 0.8;
    margin-bottom: 16px;
    opacity: 0.6;
  }}
  .quote-text {{
    font-size: 20px;
    font-style: italic;
    color: var(--white);
    line-height: 1.6;
    margin-bottom: 16px;
  }}
  .quote-attr {{
    font-size: 14px;
    color: var(--blue);
    font-weight: 600;
  }}

  /* ── Footer ── */
  footer {{
    text-align: center;
    padding: 32px;
    font-size: 12px;
    color: #8899AA;
  }}
  footer span {{ color: var(--blue); font-weight: 600; }}
</style>
</head>
<body>

<div class="hero">
  <div class="hero-label">Value Realization Story</div>
  <div class="hero-title">{customer_name}</div>
  <div class="hero-sub">{industry}</div>
  <div class="tags">
    {tags_html}
  </div>
</div>

<div class="metric-bar">
  <div class="metric">
    <div class="metric-label">Time to Value</div>
    <div class="metric-value">{time_to_value}</div>
  </div>
  <div class="metric">
    <div class="metric-label">Use Cases</div>
    <div class="metric-value">{tag_count}</div>
  </div>
  <div class="metric">
    <div class="metric-label">Outcomes Tracked</div>
    <div class="metric-value">{outcome_count}</div>
  </div>
  <div class="metric">
    <div class="metric-label">Generated</div>
    <div class="metric-value">{gen_date}</div>
  </div>
</div>

<div class="body">

  <div class="card">
    <div class="card-icon">⚡</div>
    <div class="card-title">The Challenge</div>
    <div class="card-body">{business_problem}</div>
  </div>

  <div class="card">
    <div class="card-icon">❄</div>
    <div class="card-title">The Solution</div>
    <div class="card-body">{solution_overview}</div>
  </div>

  <div class="card full">
    <div class="card-icon">📈</div>
    <div class="card-title">Business Outcomes</div>
    <ul class="outcome-list">
      {outcomes_html}
    </ul>
  </div>

  {quote_html}

</div>

<footer>
  Generated by <span>❄ Snowflake Value Realization App</span> &nbsp;·&nbsp; {gen_date}
</footer>

</body>
</html>"""


def generate_html(fields: dict) -> str:
    tags_html = "\n    ".join(
        f'<span class="tag">{t}</span>' for t in fields.get("use_case_tags", [])
    )
    outcomes_html = "\n      ".join(
        f'<li><span class="bullet"></span>{o}</li>' for o in fields.get("outcomes", [])
    )
    quote = fields.get("quote") or ""
    quote_html = ""
    if quote.strip():
        quote_html = f"""  <div class="quote-block">
    <div class="quote-mark">\u201c</div>
    <div class="quote-text">{quote}</div>
    <div class="quote-attr">\u2014 {fields.get('customer_name', '')}</div>
  </div>"""

    return HTML_TEMPLATE.format(
        customer_name=fields.get("customer_name", ""),
        industry=fields.get("industry", ""),
        tags_html=tags_html,
        time_to_value=fields.get("time_to_value", "N/A"),
        tag_count=len(fields.get("use_case_tags", [])),
        outcome_count=len(fields.get("outcomes", [])),
        gen_date=str(date.today()),
        business_problem=fields.get("business_problem", ""),
        solution_overview=fields.get("solution_overview", ""),
        outcomes_html=outcomes_html,
        quote_html=quote_html,
    )


# ═══════════════════════════════════════════════════════════════════════════
# STREAMLIT UI
# ═══════════════════════════════════════════════════════════════════════════

def _sf_blue(text: str) -> str:
    return f'<span style="color:#29B5E8;font-weight:700">{text}</span>'


def _step_bar(current: int):
    """Render a 3-step progress bar. current = 1, 2, or 3."""
    steps = [
        ("1", "Describe", "Paste your use case text"),
        ("2", "Review",   "AI extracts & you confirm"),
        ("3", "Download", "PPTX · README · HTML"),
    ]
    parts = []
    for i, (num, label, sub) in enumerate(steps, 1):
        if i < current:
            circle_style = "background:#29B5E8;color:white;"
            label_style  = "color:#29B5E8;font-weight:700;"
            sub_style    = "color:#8BB8D8;"
            check = "✓"
        elif i == current:
            circle_style = "background:#0A2F4E;color:#29B5E8;border:2px solid #29B5E8;"
            label_style  = "color:#0A2F4E;font-weight:700;"
            sub_style    = "color:#555;"
            check = num
        else:
            circle_style = "background:#DDE4EC;color:#999;"
            label_style  = "color:#AAA;"
            sub_style    = "color:#CCC;"
            check = num

        connector = (
            '<div style="flex:1;height:2px;background:#DDE4EC;margin:0 8px;align-self:center;'
            'margin-top:-24px;"></div>'
            if i < 3 else ""
        )
        parts.append(f"""
          <div style="display:flex;flex-direction:column;align-items:center;min-width:120px;">
            <div style="width:36px;height:36px;border-radius:50%;{circle_style}
                        display:flex;align-items:center;justify-content:center;
                        font-weight:700;font-size:15px;margin-bottom:6px;">{check}</div>
            <div style="font-size:13px;{label_style}">{label}</div>
            <div style="font-size:11px;{sub_style}text-align:center;max-width:110px;">{sub}</div>
          </div>
          {connector}
        """)

    html = f"""
    <div style="display:flex;align-items:flex-start;justify-content:center;
                background:white;border-radius:12px;padding:20px 32px;
                margin-bottom:24px;box-shadow:0 1px 6px rgba(10,47,78,0.08);">
      {''.join(parts)}
    </div>
    """
    st.markdown(html, unsafe_allow_html=True)


def main():
    st.set_page_config(
        page_title="Value Realization App",
        page_icon="❄",
        layout="wide",
    )

    # ── Global CSS ───────────────────────────────────────────────────────────
    st.markdown("""
    <style>
      [data-testid="stAppViewContainer"] {background:#F4F7FA;}
      .stButton>button {
        background:#29B5E8;color:white;border:none;border-radius:8px;
        padding:10px 24px;font-weight:600;font-size:15px;
      }
      .stButton>button:hover {background:#00A1E0;}
      .stTextArea textarea {font-size:14px;}
      div[data-testid="stExpander"] details {border:none;}
      .output-header {
        background:linear-gradient(135deg,#0A2F4E,#1a4d7a);
        border-radius:10px;padding:20px 28px;margin-bottom:16px;color:white;
      }
      .field-help {font-size:12px;color:#7890A0;margin-top:2px;margin-bottom:8px;}
      .quality-warn {background:#FFF8E1;border-left:4px solid #FFB300;
                     border-radius:6px;padding:14px 18px;margin-bottom:16px;}
      .quality-ok   {background:#E8F5E9;border-left:4px solid #4CAF50;
                     border-radius:6px;padding:14px 18px;margin-bottom:16px;}
    </style>
    """, unsafe_allow_html=True)

    # ── Page header ──────────────────────────────────────────────────────────
    st.markdown("""
    <div style="background:linear-gradient(135deg,#0A2F4E 0%,#1a4d7a 100%);
                border-radius:14px;padding:32px 40px;margin-bottom:20px;">
      <p style="font-size:11px;letter-spacing:3px;color:#29B5E8;
                font-weight:700;text-transform:uppercase;margin:0 0 8px">
        ❄ Snowflake</p>
      <h1 style="color:white;font-size:34px;font-weight:800;margin:0 0 8px">
        Value Realization App</h1>
      <p style="color:#8BB8D8;font-size:15px;margin:0">
        Turn a raw use case description into a branded PPTX deck, a GitHub README,
        and an HTML portfolio card — powered by Cortex AI.</p>
    </div>
    """, unsafe_allow_html=True)

    # ── Session state init ───────────────────────────────────────────────────
    if "fields"      not in st.session_state: st.session_state.fields      = None
    if "extraction"  not in st.session_state: st.session_state.extraction  = False
    if "pptx_bytes"  not in st.session_state: st.session_state.pptx_bytes  = None
    if "html_str"    not in st.session_state: st.session_state.html_str    = None

    # Determine current step for progress bar
    if st.session_state.pptx_bytes:
        _step_bar(3)
    elif st.session_state.extraction:
        _step_bar(2)
    else:
        _step_bar(1)

    # ════════════════════════════════════════════════════════════════════════
    # STEP 1 — INPUT
    # ════════════════════════════════════════════════════════════════════════
    st.markdown("### Step 1 — Describe the Use Case")

    # What-you-get row
    st.markdown("""
    <div style="display:flex;gap:12px;margin-bottom:20px;">
      <div style="flex:1;background:white;border-radius:10px;padding:16px 18px;
                  border-top:3px solid #29B5E8;box-shadow:0 1px 4px rgba(10,47,78,0.07);">
        <div style="font-size:20px;margin-bottom:6px;">📊</div>
        <div style="font-weight:700;font-size:13px;color:#0A2F4E;">PowerPoint Deck</div>
        <div style="font-size:12px;color:#7890A0;margin-top:4px;">
          6-slide Snowflake-branded story — title, challenge, solution, outcomes, why Snowflake, testimonial.</div>
      </div>
      <div style="flex:1;background:white;border-radius:10px;padding:16px 18px;
                  border-top:3px solid #29B5E8;box-shadow:0 1px 4px rgba(10,47,78,0.07);">
        <div style="font-size:20px;margin-bottom:6px;">🐙</div>
        <div style="font-weight:700;font-size:13px;color:#0A2F4E;">GitHub README</div>
        <div style="font-size:12px;color:#7890A0;margin-top:4px;">
          Markdown file pushed directly to a GitHub repo via the API — ready for your portfolio or team wiki.</div>
      </div>
      <div style="flex:1;background:white;border-radius:10px;padding:16px 18px;
                  border-top:3px solid #29B5E8;box-shadow:0 1px 4px rgba(10,47,78,0.07);">
        <div style="font-size:20px;margin-bottom:6px;">🌐</div>
        <div style="font-weight:700;font-size:13px;color:#0A2F4E;">HTML Portfolio Card</div>
        <div style="font-size:12px;color:#7890A0;margin-top:4px;">
          Self-contained .html file with Snowflake branding — embed in a site, share via email, or host anywhere.</div>
      </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("""
    <div style="background:#E3F2FD;border-left:4px solid #29B5E8;border-radius:6px;
                padding:12px 16px;margin-bottom:12px;font-size:13px;color:#1a3a52;">
      <strong>What to paste here:</strong> Cortex AI needs <em>prose</em> — a customer story, email thread,
      discovery notes, or case study bullet points. The more detail you include (company name, industry,
      specific metrics, Snowflake products used, timeline), the richer the outputs will be.
      <br><br>
      <strong>What not to paste:</strong> URLs, slide titles, or single-sentence summaries will produce
      "Unknown" fields because there's not enough context for the AI to work with.
    </div>
    """, unsafe_allow_html=True)

    description = st.text_area(
        label="Use case description",
        placeholder=(
            "Example:\n\n"
            "Acme Corp, a Fortune 500 retailer, was struggling with data silos across 12 "
            "business units — each team ran their own warehouse, causing 6-hour reporting "
            "delays and $4M in duplicated infrastructure costs. After migrating to the "
            "Snowflake Data Cloud, they consolidated all data sources into a single platform, "
            "reduced query times by 80%, cut cloud costs by $2M annually, and launched a "
            "real-time inventory recommendation engine powered by Cortex AI — all within "
            "4 months. The VP of Data Engineering said: 'We went from 12 silos to one "
            "source of truth overnight.'"
        ),
        height=220,
        label_visibility="collapsed",
    )

    extract_btn = st.button("✦  Extract Fields with Cortex AI", use_container_width=False)

    if extract_btn and description.strip():
        with st.spinner("Cortex AI is reading your description and extracting structured fields…"):
            try:
                fields = extract_fields(description)
                st.session_state.fields     = fields
                st.session_state.extraction = True
                st.session_state.pptx_bytes = None
                st.session_state.html_str   = None
                st.rerun()
            except Exception as e:
                st.error(f"Extraction failed: {e}")

    elif extract_btn:
        st.warning("Please enter a use case description first.")

    # ════════════════════════════════════════════════════════════════════════
    # STEP 2 — REVIEW & EDIT
    # ════════════════════════════════════════════════════════════════════════
    if st.session_state.extraction and st.session_state.fields:
        f = st.session_state.fields
        st.markdown("---")
        st.markdown("### Step 2 — Review & Edit Extracted Fields")

        # Quality check — count Unknown / empty fields
        key_fields = [
            f.get("customer_name", ""), f.get("industry", ""),
            f.get("business_problem", ""), f.get("solution_overview", ""),
        ]
        unknown_count = sum(1 for v in key_fields if not v or v.strip().lower() == "unknown")
        outcomes_count = len(f.get("outcomes", []))

        if unknown_count >= 3:
            st.markdown(f"""
            <div class="quality-warn">
              <strong>⚠ Most fields came back as "Unknown"</strong> — this usually means the input
              didn't contain enough detail for the AI to work with (e.g. a URL was pasted instead
              of a description). <br><br>
              You can <strong>edit the fields below manually</strong>, or scroll back up and
              re-enter a richer description with customer name, industry, problem, solution, and metrics.
            </div>
            """, unsafe_allow_html=True)
        elif unknown_count >= 1 or outcomes_count == 0:
            st.markdown(f"""
            <div class="quality-warn">
              <strong>Some fields need attention</strong> — {unknown_count} key field(s) are "Unknown"
              and {outcomes_count} business outcome(s) were found. Fill in the blanks below before generating.
            </div>
            """, unsafe_allow_html=True)
        else:
            st.markdown(f"""
            <div class="quality-ok">
              <strong>✓ Extraction looks good</strong> — all key fields were found and
              {outcomes_count} business outcome(s) extracted. Review the details below and confirm.
            </div>
            """, unsafe_allow_html=True)

        st.markdown("""
        <p style="font-size:13px;color:#556677;margin-bottom:16px;">
          Cortex AI extracted these fields from your description.
          <strong>Review each one and correct anything that looks wrong</strong> before generating outputs —
          these values drive every slide, README section, and HTML card.
        </p>
        """, unsafe_allow_html=True)

        with st.form("fields_form"):
            col1, col2 = st.columns(2)
            with col1:
                cname = st.text_input(
                    "Customer Name",
                    value=f.get("customer_name", ""),
                    help="Used as the headline on every output. Use the full company name.",
                )
                industry = st.text_input(
                    "Industry",
                    value=f.get("industry", ""),
                    help="e.g. Financial Services, Retail, Healthcare, Manufacturing",
                )
                ttv = st.text_input(
                    "Time to Value",
                    value=f.get("time_to_value", ""),
                    help="How long until the customer saw results — e.g. '3 months', '6 weeks'",
                )
            with col2:
                all_tags = [
                    "Data Engineering", "AI/ML", "Data Sharing", "Cost Optimization",
                    "Governance", "Real-Time Analytics", "Application Development", "Migration",
                ]
                tags = st.multiselect(
                    "Use Case Tags",
                    options=all_tags,
                    default=[t for t in f.get("use_case_tags", []) if t in all_tags],
                    help="Select all Snowflake capability areas this use case covers",
                )

            st.markdown("**Business Problem**")
            st.markdown('<p class="field-help">2–3 sentences describing what the customer was struggling with before Snowflake. Appears on Slide 2 and in the Challenge section.</p>', unsafe_allow_html=True)
            problem = st.text_area(
                label="Business Problem", value=f.get("business_problem", ""),
                height=100, label_visibility="collapsed",
            )

            st.markdown("**Solution Overview**")
            st.markdown('<p class="field-help">What Snowflake products, features, or approaches were used. Appears on Slide 3 and in the Solution section.</p>', unsafe_allow_html=True)
            solution = st.text_area(
                label="Solution Overview", value=f.get("solution_overview", ""),
                height=100, label_visibility="collapsed",
            )

            st.markdown("**Business Outcomes** — one per line")
            st.markdown('<p class="field-help">Quantified results wherever possible (e.g. "80% faster queries", "$2M cost savings"). Each line becomes a card on Slide 4 and a bullet in the README.</p>', unsafe_allow_html=True)
            outcomes_raw = st.text_area(
                label="outcomes",
                value="\n".join(f.get("outcomes", [])),
                height=130,
                label_visibility="collapsed",
            )

            st.markdown("**Quote / Testimonial** (optional)")
            st.markdown('<p class="field-help">A direct quote from a stakeholder. If present, adds a testimonial slide (Slide 6) and a highlighted quote block in the HTML card.</p>', unsafe_allow_html=True)
            quote = st.text_area(
                label="Quote", value=f.get("quote") or "",
                height=80, label_visibility="collapsed",
            )

            confirm_btn = st.form_submit_button(
                "✓  Confirm & Generate All Outputs", use_container_width=True
            )

        if confirm_btn:
            updated = {
                "customer_name":    cname,
                "industry":         industry,
                "time_to_value":    ttv,
                "use_case_tags":    tags,
                "business_problem": problem,
                "solution_overview":solution,
                "quote":            quote if quote.strip() else None,
                "outcomes":         [o.strip() for o in outcomes_raw.splitlines() if o.strip()],
            }
            st.session_state.fields = updated
            with st.spinner("Generating PPTX and HTML…"):
                st.session_state.pptx_bytes = generate_pptx(updated)
                st.session_state.html_str   = generate_html(updated)
            st.rerun()

    # ════════════════════════════════════════════════════════════════════════
    # STEP 3 — OUTPUTS
    # ════════════════════════════════════════════════════════════════════════
    if st.session_state.pptx_bytes and st.session_state.html_str:
        f     = st.session_state.fields
        cname = f.get("customer_name", "customer").replace(" ", "_")

        st.markdown("---")
        col_hdr, col_reset = st.columns([5, 1])
        with col_hdr:
            st.markdown("### Step 3 — Download Your Assets")
        with col_reset:
            if st.button("↺  Start Over", help="Clear all fields and start a new use case"):
                for key in ["fields", "extraction", "pptx_bytes", "html_str"]:
                    st.session_state[key] = None if key != "extraction" else False
                st.rerun()

        tab_pptx, tab_github, tab_html = st.tabs(
            ["📊  PowerPoint Deck", "🐙  GitHub README", "🌐  HTML Portfolio Card"]
        )

        # ── Tab 1 — PPTX ─────────────────────────────────────────────────────
        with tab_pptx:
            st.markdown(f"""
            <div class="output-header">
              <p style="margin:0;font-size:12px;letter-spacing:2px;
                        text-transform:uppercase;color:#29B5E8">PowerPoint</p>
              <h3 style="margin:6px 0 4px;color:white">
                {f.get('customer_name')} Value Realization Story</h3>
              <p style="margin:0;color:#8BB8D8;font-size:14px">
                {len(f.get('outcomes',[]))} outcomes  ·
                {len(f.get('use_case_tags',[]))} use case tags  ·
                {'Testimonial included' if f.get('quote') else 'No testimonial'}
              </p>
            </div>
            """, unsafe_allow_html=True)

            st.download_button(
                label="⬇  Download PPTX",
                data=st.session_state.pptx_bytes,
                file_name=f"{cname}_value_realization.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

            with st.expander("Slide outline"):
                slides_desc = [
                    ("1", "Title",            f"Customer: {f.get('customer_name')}"),
                    ("2", "The Challenge",     f.get("business_problem", "")[:80] + "…"),
                    ("3", "The Solution",      f.get("solution_overview", "")[:80] + "…"),
                    ("4", "Business Outcomes", f"{len(f.get('outcomes',[]))} bullet outcomes"),
                    ("5", "Why Snowflake",     f"Time to Value: {f.get('time_to_value')}"),
                ]
                if f.get("quote"):
                    slides_desc.append(("6", "Testimonial", f.get("quote", "")[:60] + "…"))
                for num, title, detail in slides_desc:
                    st.markdown(f"**Slide {num} — {title}**: {detail}")

        # ── Tab 2 — GitHub ────────────────────────────────────────────────────
        with tab_github:
            st.markdown("""
            <div class="output-header">
              <p style="margin:0;font-size:12px;letter-spacing:2px;
                        text-transform:uppercase;color:#29B5E8">GitHub</p>
              <h3 style="margin:6px 0 4px;color:white">Push README to GitHub</h3>
              <p style="margin:0;color:#8BB8D8;font-size:14px">
                Creates or updates README.md in the target repo via the GitHub API.</p>
            </div>
            """, unsafe_allow_html=True)

            readme_content = build_readme(f)

            col_l, col_r = st.columns(2)
            with col_l:
                gh_owner = st.text_input("GitHub username or org", key="gh_owner")
                gh_repo  = st.text_input("Repository name",        key="gh_repo")
                create_new = st.checkbox("Create repo if it doesn't exist", value=True)
                private    = st.checkbox("Private repo", value=False)

            with col_r:
                st.markdown("**README preview**")
                st.code(readme_content[:800] + "\n…(truncated)", language="markdown")

            push_btn = st.button("🐙  Push README to GitHub", use_container_width=True)

            if push_btn:
                if not gh_owner or not gh_repo:
                    st.warning("Enter the GitHub owner and repo name.")
                else:
                    try:
                        import _snowflake
                        token = _snowflake.get_generic_secret_string("github_pat")
                    except Exception:
                        st.error(
                            "Could not read `github_pat` secret. "
                            "Make sure the STREAMLIT object was created with the "
                            "`SECRETS = ('github_pat' = ...)` clause."
                        )
                        st.stop()

                    with st.spinner("Talking to GitHub…"):
                        try:
                            if create_new:
                                create_resp = create_github_repo(token, gh_owner, gh_repo, private)
                                if "id" in create_resp:
                                    st.info(f"Repo created: {create_resp.get('html_url')}")
                                elif "message" in create_resp and "already exists" not in create_resp["message"]:
                                    st.warning(f"Repo creation: {create_resp['message']}")

                            result = push_readme_to_github(token, gh_owner, gh_repo, readme_content)

                            if "content" in result:
                                html_url = result["content"].get("html_url", "")
                                st.success(f"README pushed successfully! [{html_url}]({html_url})")
                            else:
                                st.error(f"GitHub API response: {result}")
                        except Exception as e:
                            st.error(f"Push failed: {e}")

            # Always offer local download too
            st.download_button(
                "⬇  Download README.md locally",
                data=readme_content,
                file_name=f"{cname}_README.md",
                mime="text/markdown",
            )

        # ── Tab 3 — HTML ──────────────────────────────────────────────────────
        with tab_html:
            st.markdown("""
            <div class="output-header">
              <p style="margin:0;font-size:12px;letter-spacing:2px;
                        text-transform:uppercase;color:#29B5E8">HTML Portfolio</p>
              <h3 style="margin:6px 0 4px;color:white">Self-Contained Portfolio Card</h3>
              <p style="margin:0;color:#8BB8D8;font-size:14px">
                Single .html file — paste into a portfolio, website, or email.</p>
            </div>
            """, unsafe_allow_html=True)

            st.download_button(
                label="⬇  Download HTML Card",
                data=st.session_state.html_str,
                file_name=f"{cname}_portfolio.html",
                mime="text/html",
                use_container_width=True,
            )

            with st.expander("Preview HTML card", expanded=True):
                import streamlit.components.v1 as components
                components.html(st.session_state.html_str, height=700, scrolling=True)


if __name__ == "__main__":
    main()
